import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Initialize Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Corporate Light Mode
    c_navy = RGBColor(10, 25, 47)       # Deep Corporate Navy
    c_gold = RGBColor(197, 168, 128)    # Premium Gold
    c_gold_dark = RGBColor(163, 131, 88) # Darker Gold for text
    c_bg_light = RGBColor(248, 249, 250) # Light grey bg
    c_white = RGBColor(255, 255, 255)
    c_text_dark = RGBColor(51, 51, 51)
    c_text_muted = RGBColor(102, 102, 102)
    c_border = RGBColor(226, 232, 240)

    # Helper function to set slide background
    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper function to create a stylized slide header
    def add_header(slide, title_text, category="PHÂN TÍCH VĨ MÔ 2026"):
        # Category Badge
        badge_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8), Inches(0.3))
        tf_b = badge_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.text = category.upper()
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = c_gold_dark

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.8))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = "Calibri"
        p_t.font.size = Pt(28)
        p_t.font.bold = True
        p_t.font.color.rgb = c_navy

        # Decorative bottom divider line under header
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = c_gold
        line.line.color.rgb = c_gold

    # Helper function to create formatted container boxes
    def add_card(slide, left, top, width, height, title, content_list, border_color=None, bg_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color if bg_color else c_white
        card.line.color.rgb = border_color if border_color else c_border
        card.line.width = Pt(1.5)

        # Content TextBox inside card
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        # Card Title
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Calibri"
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = c_navy
        p_t.space_after = Pt(12)

        # Bullet Items
        for idx, item in enumerate(content_list):
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.name = "Calibri"
            p.font.size = Pt(13)
            p.font.color.rgb = c_text_dark
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Trang bìa)
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide1 = prs.slides.add_slide(slide_layout)
    set_bg(slide1, c_navy)

    # Elegant gold line on the left side
    gold_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0), Inches(0.15), Inches(7.5))
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = c_gold
    gold_bar.line.fill.background()

    # Title & Subtitle box
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(3))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "BÁO CÁO CHIẾN LƯỢC"
    p1.font.name = "Calibri"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = c_gold
    p1.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "BỐI CẢNH KINH TẾ - CHÍNH TRỊ TOÀN CẦU 2026"
    p2.font.name = "Calibri"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = c_white
    p2.space_after = Pt(14)

    p3 = tf1.add_paragraph()
    p3.text = "Đánh giá Tác động vĩ mô & Khuyến nghị Hành động cho Doanh nghiệp Việt Nam"
    p3.font.name = "Calibri"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(200, 200, 200)

    # Footer/Metadata
    meta_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(6), Inches(1))
    tf_meta = meta_box.text_frame
    pm = tf_meta.paragraphs[0]
    pm.text = "Thực hiện bởi: Nhóm chuyên gia Antigravity\nNgày báo cáo: 29/05/2026"
    pm.font.name = "Calibri"
    pm.font.size = Pt(11)
    pm.font.color.rgb = c_gold

    # ----------------------------------------------------
    # SLIDE 2: Executive Summary (Tóm tắt vĩ mô)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_bg(slide2, c_bg_light)
    add_header(slide2, "TÓM TẮT ĐIỀU HÀNH: BỨC TRANH TOÀN CẦU NĂM 2026")

    # 4 Columns for 4 Main forces
    add_card(slide2, Inches(0.6), Inches(1.8), Inches(2.8), Inches(4.8), 
             "1. Thuế quan Mỹ", 
             ["Mỹ áp thuế bảo hộ đồng đều 10% toàn cầu.", "Lạm phát Mỹ neo cao 3,8% do chi phí đẩy.", "Fed giữ lãi suất cao 3,5%-3,75% trì hoãn giảm."])

    add_card(slide2, Inches(3.6), Inches(1.8), Inches(2.8), Inches(4.8), 
             "2. Bảo hộ Xanh EU", 
             ["CBAM bước vào giai đoạn tài chính bắt buộc 1/1/2026.", "Hạn chót thực thi EUDR về phá rừng (30/12/2026).", "Áp 100% thuế khí thải hàng hải EU ETS."])

    add_card(slide2, Inches(6.6), Inches(1.8), Inches(2.8), Inches(4.8), 
             "3. FDI & Chuỗi SEA", 
             ["FDI đổ dồn vào ASEAN đạt kỷ lục >200 tỷ USD/năm.", "Malaysia đột phá bán dẫn, Thái Lan thắt chặt EV 3.5.", "DEFA số hóa và ASW giảm 60% thông quan."])

    add_card(slide2, Inches(9.6), Inches(1.8), Inches(3.1), Inches(4.8), 
             "4. Logistics Trung Đông", 
             ["Xung đột Mỹ-Iran phong tỏa eo biển Hormuz.", "Đi vòng qua Mũi Hảo Vọng là tiêu chuẩn bình thường.", "Cước biển tăng 2-3 lần, dầu thô vượt $120/thùng."])

    # ----------------------------------------------------
    # SLIDE 3: USA Region
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_bg(slide3, c_bg_light)
    add_header(slide3, "KHU VỰC MỸ: THUẾ QUAN BẢO HỘ & CƠ HỘI BÁN DẪN", "ĐÁNH GIÁ CHI TIẾT KHU VỰC")

    # Impact badge
    impact_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.4), Inches(2.2), Inches(0.6))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(220, 53, 69) # Red color for high impact
    impact_box.line.color.rgb = RGBColor(220, 53, 69)
    p_imp = impact_box.text_frame.paragraphs[0]
    p_imp.alignment = PP_ALIGN.CENTER
    p_imp.text = "TÁC ĐỘNG: CAO"
    p_imp.font.name = "Calibri"
    p_imp.font.size = Pt(14)
    p_imp.font.bold = True
    p_imp.font.color.rgb = c_white

    add_card(slide3, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), 
             "Thách thức & Rủi ro ngắn hạn", 
             ["Áp thuế quan đồng đều 10% theo Mục 122 Đạo luật Thương mại.", "Đe dọa trực tiếp biên lợi nhuận của dệt may, da giày xuất khẩu.", "Siết chặt UFLPA chống lao động cưỡng bức, nguy cơ áp thuế phạt chuyển tải 40% nếu lẩn tránh nguồn gốc.", "Sức mua người tiêu dùng Mỹ suy yếu do lãi suất Fed cao 3,5%-3,75%."])

    add_card(slide3, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8), 
             "Cơ hội & Đòn bẩy chiến lược", 
             ["Căng thẳng Mỹ - Trung thúc đẩy làn sóng FDI dịch chuyển công nghệ bán dẫn và điện tử thế hệ mới sang Việt Nam.", "Luật Công nghiệp Công nghệ số Việt Nam (1/1/2026) ưu đãi thuế kỷ lục 5% trong suốt 37 năm.", "Việt Nam bùng nổ thu hút siêu dự án bán dẫn (Samsung 4 tỷ USD ở Thái Nguyên, Viettel bán dẫn tự chủ)."])

    # ----------------------------------------------------
    # SLIDE 4: EU Region
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_bg(slide4, c_bg_light)
    add_header(slide4, "KHU VỰC EU: TIÊU CHUẨN XANH PHÁP LÝ BẮT BUỘC", "ĐÁNH GIÁ CHI TIẾT KHU VỰC")

    # Impact badge
    impact_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.4), Inches(2.2), Inches(0.6))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(220, 53, 69)
    impact_box.line.color.rgb = RGBColor(220, 53, 69)
    p_imp = impact_box.text_frame.paragraphs[0]
    p_imp.alignment = PP_ALIGN.CENTER
    p_imp.text = "TÁC ĐỘNG: CAO"
    p_imp.font.name = "Calibri"
    p_imp.font.size = Pt(14)
    p_imp.font.bold = True
    p_imp.font.color.rgb = c_white

    add_card(slide4, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), 
             "Rủi ro tuân thủ & Rào cản kỹ thuật", 
             ["CBAM Definitive Phase có hiệu lực hoàn toàn từ 1/1/2026.", "Hạn chót thực thi EUDR bắt buộc vào 30/12/2026 đối với nông sản, rủi ro đứt gãy xuất khẩu nếu thiếu định vị GPS.", "Thẻ vàng IUU thủy sản chưa được gỡ sau đợt kiểm tra thứ 5 (3/2026) gây chậm trễ hải quan và tăng chi phí.", "Cấm tiêu hủy hàng dệt may chưa bán (từ 19/7/2026), áp dụng ESPR và hộ chiếu số DPP."])

    add_card(slide4, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8), 
             "Cơ hội & Lợi thế đi đầu", 
             ["Nâng cấp quan hệ Đối tác Chiến lược Toàn diện VN-EU đầu năm 2026, thu hút FDI xanh dồi dào.", "Lợi thế vượt trội từ thuế suất 0% của EVFTA (năm thứ 6) so với các đối thủ ASEAN.", "Gói cải cách Omnibus I lùi CSDDD/CSRD đến 2028-2029 và nới lỏng ngưỡng áp dụng, giúp các SME Việt Nam dễ thở hơn."])

    # ----------------------------------------------------
    # SLIDE 5: China Region
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_bg(slide5, c_bg_light)
    add_header(slide5, "KHU VỰC TRUNG QUỐC: ĐỊNH HÌNH LẠI CHUỖI CUNG ỨNG THƯỢNG NGUỒN", "ĐÁNH GIÁ CHI TIẾT KHU VỰC")

    # Impact badge
    impact_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.4), Inches(2.2), Inches(0.6))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(255, 193, 7) # Yellow for medium impact
    impact_box.line.color.rgb = RGBColor(255, 193, 7)
    p_imp = impact_box.text_frame.paragraphs[0]
    p_imp.alignment = PP_ALIGN.CENTER
    p_imp.text = "TÁC ĐỘNG: TRUNG BÌNH"
    p_imp.font.name = "Calibri"
    p_imp.font.size = Pt(14)
    p_imp.font.bold = True
    p_imp.font.color.rgb = c_navy

    add_card(slide5, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), 
             "Thách thức & Cạnh tranh gay gắt", 
             ["Trung Quốc xả hàng tồn kho giá rẻ qua các sàn TMĐT (Temu, Shein) gây sức ép khốc liệt lên doanh nghiệp nội địa.", "GACC siết chặt tiêu chuẩn chất lượng (SPS) đối với xuất khẩu nông sản tiểu ngạch của Việt Nam.", "Đồng RMB mạnh làm tăng chi phí nhập khẩu nguyên vật liệu đầu vào của doanh nghiệp Việt.", "Cạnh tranh thu hút FDI bán dẫn với các gói hỗ trợ khủng của Bắc Kinh."])

    add_card(slide5, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8), 
             "Cơ hội hợp tác vĩ mô", 
             ["Tuyến đường sắt liên vận khổ tiêu chuẩn (1.435mm) Lào Cai - Hà Nội - Hải Phòng vận hành giúp giảm 30%-40% chi phí logistics.", "Tiếp cận nguồn nguyên liệu xanh, đạt chuẩn công nghệ tiên tiến của Trung Quốc để xuất khẩu đi EU/Mỹ.", "Cơ hội thu hút các dự án pin công nghệ cao khi Bắc Kinh giảm hoàn thuế xuất khẩu pin xuống 6% từ 4/2026."])

    # ----------------------------------------------------
    # SLIDE 6: Southeast Asia Region
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_bg(slide6, c_bg_light)
    add_header(slide6, "ĐÔNG NAM Á: CUỘC ĐUA FDI KHỐC LIỆT & LIÊN KẾT SỐ NỘI KHỐI", "ĐÁNH GIÁ CHI TIẾT KHU VỰC")

    # Impact badge
    impact_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.4), Inches(2.2), Inches(0.6))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(220, 53, 69)
    impact_box.line.color.rgb = RGBColor(220, 53, 69)
    p_imp = impact_box.text_frame.paragraphs[0]
    p_imp.alignment = PP_ALIGN.CENTER
    p_imp.text = "TÁC ĐỘNG: CAO"
    p_imp.font.name = "Calibri"
    p_imp.font.size = Pt(14)
    p_imp.font.bold = True
    p_imp.font.color.rgb = c_white

    add_card(slide6, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), 
             "Cạnh tranh bứt phá của đối thủ", 
             ["Malaysia vượt trội trong thu hút FDI bán dẫn chất lượng cao nhờ Chiến lược Bán dẫn Quốc gia (NSS): Gói R&D 92 triệu RM đóng gói tiên tiến chiplet, quỹ 550 triệu RM phát triển hệ sinh thái và quỹ tín dụng ưu đãi 500 triệu RM.", "Thái Lan thắt chặt chính sách EV 3.5 yêu cầu tỷ lệ lắp ráp nội địa 1:2 năm 2026, ưu đãi tiêu thụ đặc biệt BEV chỉ 2%.", "Cạnh tranh lao động giá rẻ từ Campuchia, Myanmar."])

    add_card(slide6, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8), 
             "Động lực liên kết vùng", 
             ["Hiệp định Khung Kinh tế Số ASEAN (DEFA) đặt mục tiêu đưa kinh tế số đạt 2.000 tỷ USD vào năm 2030.", "Cơ chế một cửa ASEAN (ASW) tích hợp AI/Blockchain rút ngắn 60% thời gian thông quan nội khối, giúp hàng nông sản tươi sống của Việt Nam xuất sang Singapore/Malaysia thuận tiện.", "Cơ hội liên kết chuỗi đóng gói nâng cấp với Malaysia và R&D với Singapore."])

    # ----------------------------------------------------
    # SLIDE 7: Middle East Region
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_bg(slide7, c_bg_light)
    add_header(slide7, "TRUNG ĐÔNG: KHỦNG HOẢNG LOGISTICS & KHUNG Halal VÀNG", "ĐÁNH GIÁ CHI TIẾT KHU VỰC")

    # Impact badge
    impact_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.4), Inches(2.2), Inches(0.6))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(220, 53, 69)
    impact_box.line.color.rgb = RGBColor(220, 53, 69)
    p_imp = impact_box.text_frame.paragraphs[0]
    p_imp.alignment = PP_ALIGN.CENTER
    p_imp.text = "TÁC ĐỘNG: CAO"
    p_imp.font.name = "Calibri"
    p_imp.font.size = Pt(14)
    p_imp.font.bold = True
    p_imp.font.color.rgb = c_white

    add_card(slide7, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), 
             "Cú sốc Logistics đường biển", 
             ["Xung đột leo thang Mỹ - Iran, phong tỏa eo biển Hormuz khiến vận tải vòng qua Mũi Hảo Vọng kéo dài thêm 10-20 ngày hải trình.", "Cước vận tải biển tăng vọt 2-3 lần do phụ phí nhiên liệu và bảo hiểm tăng vọt, thiếu vỏ container rỗng tại Việt Nam.", "Giá dầu Brent vượt $120/thùng gây áp lực lạm phát chi phí đẩy toàn cầu.", "Thời gian giao hàng kéo dài làm giảm biên lợi nhuận điện tử/nông sản."])

    add_card(slide7, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8), 
             "Cửa ngõ vàng Trung Đông", 
             ["Hiệp định CEPA Việt Nam - UAE (2/2026) chính thức có hiệu lực, UAE xóa bỏ đến 99% dòng thuế cho hàng hóa Việt Nam.", "Nghị định 127/2026/NĐ-CP (4/2026) định hình khung pháp lý hoàn chỉnh cho Halal Việt Nam theo hướng 'Halal Xanh'.", "Khai phá thị trường tiêu thụ Hồi giáo GCC cực kỳ màu mỡ thông qua UAE.", "Đẩy mạnh tuyến logistics đường bộ - hàng không Sea-Air qua Dubai."])

    # ----------------------------------------------------
    # SLIDE 8: Comparison Table (Bảng so sánh)
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_bg(slide8, c_bg_light)
    add_header(slide8, "SO SÁNH BỐI CẢNH 5 KHU VỰC ĐỊA CHÍNH TRỊ NĂM 2026", "TỔNG HỢP SO SÁNH")

    # Add Table
    rows = 6
    cols = 5
    left = Inches(0.6)
    top = Inches(1.8)
    width = Inches(12.133)
    height = Inches(4.8)

    table_shape = slide8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set Column Widths
    table.columns[0].width = Inches(1.8) # Khu vực
    table.columns[1].width = Inches(3.2) # Xu hướng & Rủi ro
    table.columns[2].width = Inches(3.2) # Cơ hội nổi bật
    table.columns[3].width = Inches(1.8) # Tác động
    table.columns[4].width = Inches(2.133) # Lưu ý cốt lõi

    # Headers
    headers = ["Khu vực", "Xu hướng & Rủi ro chính", "Cơ hội lớn nhất", "Tác động", "Lưu ý cho DN Việt"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_navy
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.name = "Calibri"
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = c_white

    # Data
    data = [
        ["Mỹ (USA)", "Thuế quan 10%, UFLPA siết chặt, Fed lãi suất 3,5%-3,75%", "Dịch chuyển chuỗi bán dẫn công nghệ cao, Luật CN Công nghệ số", "CAO", "Tránh thuế phạt 40% xuất xứ, chuyển từ CMT sang ODM"],
        ["EU", "CBAM (1/1/26), EUDR (30/12/26), cước carbon hàng hải tăng vọt", "EVFTA 0% thuế quan, nới lỏng CSRD/CSDDD sang 2028-2029", "CAO", "Định vị vùng trồng GPS cho nông sản, lấy chứng chỉ xanh"],
        ["Trung Quốc", "Hàng giá rẻ qua TMĐT nội địa, SPS nông sản tiểu ngạch siết chặt", "Nguyên liệu xanh đầu vào, vận tải đường sắt liên vận khổ 1.435mm", "TRUNG BÌNH", "Liên doanh kho thông minh AI, chuyển đổi xuất khẩu chính ngạch"],
        ["Đông Nam Á", "Malaysia NSS đóng gói chip, Thái Lan EV 3.5, đua FDI gắt", "Hiệp định kinh tế số DEFA, một cửa ASW giảm 60% thông quan", "CAO", "Thiết lập liên kết đóng gói chip với Malaysia, số hóa ASW"],
        ["Trung Đông", "Xung đột Mỹ-Iran, Hormuz, đi vòng Hảo Vọng cước biển gấp 2-3", "CEPA VN-UAE xóa 99% thuế, NĐ 127/2026/NĐ-CP khung Halal Việt", "CAO", "Chuyển thương mại từ CIF sang FOB để đẩy chi phí cước biển"]
    ]

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_value
            cell.fill.solid()
            
            # Subtle alternating row colors
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = c_white
            else:
                cell.fill.fore_color.rgb = RGBColor(240, 242, 245)
                
            for paragraph in cell.text_frame.paragraphs:
                # Alignments
                if col_idx in [0, 3]:
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.alignment = PP_ALIGN.LEFT
                    
                paragraph.font.name = "Calibri"
                paragraph.font.size = Pt(11)
                
                # Colors
                if col_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = c_navy
                elif col_idx == 3:
                    paragraph.font.bold = True
                    if cell_value == "CAO":
                        paragraph.font.color.rgb = RGBColor(220, 53, 69)
                    else:
                        paragraph.font.color.rgb = c_gold_dark
                else:
                    paragraph.font.color.rgb = c_text_dark

    # ----------------------------------------------------
    # SLIDE 9: Recommendations (Khuyến nghị ngành hàng)
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(slide_layout)
    set_bg(slide9, c_bg_light)
    add_header(slide9, "KHUYẾN NGHỊ HÀNH ĐỘNG CHIẾN LƯỢC THEO NGÀNH HÀNG", "CHIẾN LƯỢC THÍCH ỨNG")

    # 3 Columns for 3 Target sectors
    add_card(slide9, Inches(0.6), Inches(1.8), Inches(3.9), Inches(4.8), 
             "1. Nhóm ngành Xuất khẩu", 
             ["**Minh bạch chuỗi cung ứng:** Rà soát nguyên liệu đầu vào phi Trung Quốc để tránh bộ lọc UFLPA Mỹ.", "**Số hóa GPS vùng trồng:** Thực hiện định vị tọa độ GPS vùng trồng đáp ứng chuẩn EUDR trước 30/12/2026.", "**Khai thác thị trường Halal:** Tận dụng CEPA Việt Nam - UAE và tiêu chuẩn Halal Xanh (Nghị định 127/2026/NĐ-CP)."])

    add_card(slide9, Inches(4.7), Inches(1.8), Inches(3.9), Inches(4.8), 
             "2. Điện tử & Bán dẫn", 
             ["**Tối ưu hóa Luật mới:** Đăng ký ưu đãi thuế thu nhập doanh nghiệp 5% trong 37 năm của Luật CN Công nghệ số.", "**Xanh hóa nhà máy:** Đầu tư điện mặt trời áp mái, giảm phát thải Scope 1 & 2 đáp ứng chuẩn chuỗi cung ứng xanh toàn cầu.", "**Liên kết nhân lực:** Hợp tác trường đại học giải quyết bài toán thiếu hụt nhân sự chất lượng cao và ngăn chảy máu chất xám."])

    add_card(slide9, Inches(8.8), Inches(1.8), Inches(3.9), Inches(4.8), 
             "3. Logistics & Chuỗi cung ứng", 
             ["**Chuyển đổi điều khoản thương mại:** Đàm phán chuyển hợp đồng CIF sang FOB để tránh phụ phí cước biển.", "**Đa dạng hóa lộ trình:** Kết hợp Sea-Air trung chuyển qua Dubai và tận dụng đường sắt khổ 1.435mm liên vận Á-Âu.", "**Số hóa hạ tầng số:** Liên doanh kho thông minh AI, ứng dụng Blockchain tối ưu hóa xoay vòng container."])

    # ----------------------------------------------------
    # SLIDE 10: Conclusion Slide (Trang kết)
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(slide_layout)
    set_bg(slide10, c_navy)

    # Gold line on top
    gold_bar_top = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    gold_bar_top.fill.solid()
    gold_bar_top.fill.fore_color.rgb = c_gold
    gold_bar_top.line.fill.background()

    # Message
    msg_box = slide10.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.5))
    tf10 = msg_box.text_frame
    tf10.word_wrap = True
    
    p10 = tf10.paragraphs[0]
    p10.alignment = PP_ALIGN.CENTER
    p10.text = "KẾT LUẬN CHIẾN LƯỢC"
    p10.font.name = "Calibri"
    p10.font.size = Pt(22)
    p10.font.bold = True
    p10.font.color.rgb = c_gold
    p10.space_after = Pt(20)

    p11 = tf10.add_paragraph()
    p11.alignment = PP_ALIGN.CENTER
    p11.text = "\"Khủng hoảng logistics là phép thử cho năng lực thích ứng nhanh,\nvà rào cản xanh chính là động lực để bứt phá chất lượng.\""
    p11.font.name = "Calibri"
    p11.font.size = Pt(26)
    p11.font.bold = True
    p11.font.color.rgb = c_white
    p11.space_after = Pt(24)

    p12 = tf10.add_paragraph()
    p12.alignment = PP_ALIGN.CENTER
    p12.text = "Chủ động Số hóa - Nâng cấp Xanh - Đa dạng hóa Chuỗi cung ứng là chìa khóa thành công."
    p12.font.name = "Calibri"
    p12.font.size = Pt(16)
    p12.font.color.rgb = RGBColor(200, 200, 200)

    # Save presentation
    output_dir = "e:/AI_STUDY/ANTIGRAVITY/output"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    prs.save(os.path.join(output_dir, "Global_Economic_Political_Summary_2026.pptx"))
    print("PowerPoint presentation generated successfully!")

if __name__ == "__main__":
    create_deck()
